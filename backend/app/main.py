from __future__ import annotations

from io import BytesIO
from pathlib import Path

from docx import Document
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from openpyxl import load_workbook
from pptx import Presentation

app = FastAPI(
    title="MAGHRABI Office Studio API",
    version="0.1.0",
    description="Document analysis and formatting engine for DOCX, XLSX and PPTX files.",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

SUPPORTED = {
    ".docx": "word",
    ".xlsx": "excel",
    ".pptx": "powerpoint",
}


@app.get("/")
def root() -> dict[str, str]:
    return {"name": "MAGHRABI Office Studio API", "status": "online", "version": "0.1.0"}


@app.get("/health")
def health() -> dict[str, str]:
    return {"status": "healthy"}


@app.get("/api/v1/formats")
def formats() -> dict[str, list[str]]:
    return {"supported": list(SUPPORTED.keys())}


def analyze_docx(data: bytes) -> dict:
    document = Document(BytesIO(data))
    paragraphs = [p for p in document.paragraphs if p.text.strip()]
    styles = sorted({p.style.name for p in paragraphs if p.style is not None})
    tables = document.tables
    headings = sum(1 for p in paragraphs if p.style and p.style.name.lower().startswith("heading"))
    return {
        "paragraphs": len(paragraphs),
        "headings": headings,
        "tables": len(tables),
        "styles": styles[:20],
    }


def analyze_xlsx(data: bytes) -> dict:
    workbook = load_workbook(BytesIO(data), read_only=True, data_only=False)
    sheets = []
    for sheet in workbook.worksheets:
        sheets.append({
            "name": sheet.title,
            "rows": sheet.max_row,
            "columns": sheet.max_column,
        })
    return {"worksheets": len(workbook.worksheets), "sheets": sheets}


def analyze_pptx(data: bytes) -> dict:
    presentation = Presentation(BytesIO(data))
    text_shapes = 0
    tables = 0
    pictures = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text_shapes += 1
            if getattr(shape, "has_table", False):
                tables += 1
            if shape.shape_type == 13:
                pictures += 1
    return {
        "slides": len(presentation.slides),
        "text_shapes": text_shapes,
        "tables": tables,
        "pictures": pictures,
    }


@app.post("/api/v1/files/analyze")
async def analyze_file(file: UploadFile = File(...)) -> dict:
    filename = file.filename or "document"
    extension = Path(filename).suffix.lower()
    document_type = SUPPORTED.get(extension)
    if not document_type:
        raise HTTPException(status_code=415, detail="Unsupported file type. Use DOCX, XLSX or PPTX.")

    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="The uploaded file is empty.")

    try:
        if extension == ".docx":
            metrics = analyze_docx(data)
        elif extension == ".xlsx":
            metrics = analyze_xlsx(data)
        else:
            metrics = analyze_pptx(data)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not parse Office file: {exc}") from exc

    return {
        "filename": filename,
        "type": document_type,
        "size_bytes": len(data),
        "status": "analyzed",
        "health_score": 70,
        "metrics": metrics,
        "issues": [
            "Formatting health rules will be expanded in V0.2",
            "Automatic repair engine is not enabled yet",
        ],
    }
